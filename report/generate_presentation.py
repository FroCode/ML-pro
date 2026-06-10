"""Generate a visually rich MedAssist Local project presentation (.pptx)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── paths ──────────────────────────────────────────────────────────────────
HERE = Path(__file__).parent
ASSETS = HERE / "assets"
OUTPUT = HERE / "MedAssist_Local_Premium_Presentation.pptx"

# ── brand palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x1D, 0x3A)
TEAL   = RGBColor(0x00, 0xB4, 0xA6)
BLUE   = RGBColor(0x4A, 0x90, 0xD9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFB)
GRAY   = RGBColor(0x5A, 0x6A, 0x7A)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
PURPLE = RGBColor(0x7B, 0x61, 0xFF)

W = Inches(13.333)
H = Inches(7.5)


# ── low-level helpers ──────────────────────────────────────────────────────
def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def _rounded(slide, l, t, w, h, fill, text="", sz=14, bold=False, tc=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = tc
    return s


def _text(slide, l, t, w, h, text, sz=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _bullets(slide, l, t, w, h, items: list[str], sz=17, color=NAVY, spacing=8):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.level = 0
    return box


def _header_bar(slide, title: str, subtitle: str = ""):
    _rect(slide, Inches(0), Inches(0), W, Inches(1.15), NAVY)
    _rect(slide, Inches(0), Inches(1.15), W, Inches(0.06), TEAL)
    _text(slide, Inches(0.55), Inches(0.18), Inches(10), Inches(0.55),
          title, sz=30, bold=True, color=WHITE)
    if subtitle:
        _text(slide, Inches(0.55), Inches(0.68), Inches(10), Inches(0.4),
              subtitle, sz=14, color=TEAL)


def _footer(slide, text: str = "MedAssist Local  |  ML Final Project  |  Track 1: Local LLM + RAG"):
    _text(slide, Inches(0.55), Inches(7.15), Inches(12), Inches(0.3),
          text, sz=9, color=GRAY, align=PP_ALIGN.RIGHT)


def _img(slide, path: Path, l, t, w, h=None):
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=w, height=h)
        return True
    return False


def _section_slide(prs, part: str, title: str, subtitle: str, img_name: str):
    slide = _blank(prs)
    if not _img(slide, ASSETS / img_name, Inches(0), Inches(0), W, H):
        _rect(slide, Inches(0), Inches(0), W, H, NAVY)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = NAVY
    overlay.fill.transparency = 0.42
    overlay.line.fill.background()

    _text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.5),
          part, sz=16, bold=True, color=TEAL)
    _text(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1.2),
          title, sz=44, bold=True, color=WHITE)
    _text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.8),
          subtitle, sz=20, color=WHITE)
    _rect(slide, Inches(0.8), Inches(5.2), Inches(2.5), Inches(0.07), TEAL)


# ── slide builders ─────────────────────────────────────────────────────────
def slide_title(prs):
    slide = _blank(prs)
    if _img(slide, ASSETS / "hero_healthcare.jpg", Inches(0), Inches(0), W, H):
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = NAVY
        overlay.fill.transparency = 0.55
        overlay.line.fill.background()

    _rounded(slide, Inches(0.6), Inches(0.5), Inches(3.2), Inches(0.45),
             TEAL, "Technical Track 1  •  Local LLM + RAG", sz=11, bold=True)

    _text(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(1.2),
          "MedAssist Local", sz=54, bold=True, color=WHITE)
    _text(slide, Inches(0.6), Inches(3.1), Inches(10), Inches(0.7),
          "Private Healthcare Q&A Assistant", sz=28, color=TEAL)
    _text(slide, Inches(0.6), Inches(4.0), Inches(10), Inches(1.2),
          "Retrieval-Augmented Generation  •  Ollama  •  ChromaDB  •  Streamlit",
          sz=16, color=WHITE)

    _rounded(slide, Inches(0.6), Inches(5.6), Inches(4.5), Inches(0.55),
             BLUE, "[Your Name(s)]  |  June 2026  |  University ML Final Lab", sz=12)

    chips = ["Healthcare", "Privacy-First", "Open Source", "Numerical Eval"]
    for i, chip in enumerate(chips):
        _rounded(slide, Inches(0.6 + i * 2.1), Inches(6.4), Inches(1.9), Inches(0.4),
                 TEAL, chip, sz=10, bold=True)


def slide_agenda(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Presentation Agenda")

    items = [
        ("01", "Problem & Motivation",        "Why healthcare needs local, cited AI answers"),
        ("02", "Technical Approach",          "RAG architecture and local LLM deployment"),
        ("03", "Dataset & Preprocessing",     "MedQuAD corpus and indexing pipeline"),
        ("04", "Implementation Plan",         "Four development phases with deliverables"),
        ("05", "Evaluation & Metrics",        "Numerical performance measurement"),
        ("06", "Demo, Timeline & Conclusion", "Live demo plan and project schedule"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.55 + col * 6.4)
        y = Inches(1.55 + row * 1.85)
        _rounded(slide, x, y, Inches(6.1), Inches(1.55), WHITE)
        _rounded(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.55), Inches(0.55),
                 TEAL, num, sz=14, bold=True)
        _text(slide, x + Inches(0.95), y + Inches(0.22), Inches(4.8), Inches(0.4),
              title, sz=17, bold=True, color=NAVY)
        _text(slide, x + Inches(0.95), y + Inches(0.65), Inches(4.8), Inches(0.6),
              desc, sz=12, color=GRAY)
    _footer(slide)


def slide_problem(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Problem & Motivation",
                "The challenge of finding trustworthy medical information")

    _img(slide, ASSETS / "doctor_consult.jpg", Inches(0.55), Inches(1.45), Inches(5.5), Inches(5.3))

    stats = [
        ("47,000+", "Medical Q&A pairs\nin MedQuAD alone"),
        ("12,000+", "NIH health topics\ncovered"),
        ("100%", "Local — data never\nleaves your machine"),
    ]
    for i, (val, lbl) in enumerate(stats):
        y = Inches(1.55 + i * 1.75)
        _rounded(slide, Inches(6.4), y, Inches(6.4), Inches(1.45), WHITE)
        _text(slide, Inches(6.6), y + Inches(0.15), Inches(1.8), Inches(0.7),
              val, sz=28, bold=True, color=TEAL)
        _text(slide, Inches(8.5), y + Inches(0.3), Inches(4.1), Inches(0.9),
              lbl, sz=13, color=GRAY)

    _bullets(slide, Inches(6.4), Inches(6.05), Inches(6.4), Inches(1.0), [
        "Cloud AI APIs expose sensitive health queries to third parties.",
        "LLMs alone can hallucinate medical facts without source grounding.",
    ], sz=13, color=NAVY)
    _footer(slide)


def slide_scope(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Project Scope & Objectives")

    objectives = [
        (TEAL,   "Application Area", "Healthcare informational Q&A\n(drug info, conditions, symptoms)"),
        (BLUE,   "Technical Track",  "Local LLM + RAG\n(Ollama + ChromaDB + Streamlit)"),
        (PURPLE, "Privacy Goal",     "100% offline inference\nNo external API calls required"),
        (ORANGE, "Evaluation",       "Numerical metrics:\nRecall@k, MRR, ROUGE-L, Latency"),
    ]
    for i, (col, title, body) in enumerate(objectives):
        x = Inches(0.55 + (i % 2) * 6.4)
        y = Inches(1.55 + (i // 2) * 2.7)
        _rounded(slide, x, y, Inches(6.1), Inches(2.35), col)
        _text(slide, x + Inches(0.3), y + Inches(0.3), Inches(5.5), Inches(0.5),
              title, sz=20, bold=True, color=WHITE)
        _text(slide, x + Inches(0.3), y + Inches(0.95), Inches(5.5), Inches(1.2),
              body, sz=15, color=WHITE)

    _rounded(slide, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.45),
             RGBColor(0xFF, 0xF3, 0xCD), "⚠  Disclaimer: Informational use only — not a substitute for professional medical advice.",
             sz=11, bold=True, tc=RGBColor(0x85, 0x60, 0x04))
    _footer(slide)


def slide_why_rag(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Why Local LLM + RAG?",
                "Combining retrieval and generation for trustworthy answers")
    _img(slide, ASSETS / "rag_vs_llm.png", Inches(0.4), Inches(1.35), Inches(12.5))
    _footer(slide)


def slide_architecture(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "System Architecture",
                "End-to-end RAG pipeline from question to cited answer")
    _img(slide, ASSETS / "architecture_diagram.png", Inches(0.3), Inches(1.3), Inches(12.7))

    steps = [
        "1. User asks via Streamlit",
        "2. Embed query (384-dim)",
        "3. Retrieve Top-5 chunks",
        "4. Build grounded prompt",
        "5. Ollama generates answer",
        "6. Return answer + citations",
    ]
    for i, s in enumerate(steps):
        _rounded(slide, Inches(0.4 + i * 2.15), Inches(6.75), Inches(2.0), Inches(0.42),
                 TEAL if i % 2 == 0 else BLUE, s, sz=8, bold=True)
    _footer(slide)


def slide_dataset(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Dataset — MedQuAD",
                "NIH-sourced medical Q&A from Hugging Face")

    _img(slide, ASSETS / "data_pipeline.png", Inches(0.4), Inches(1.35), Inches(12.5))

    details = [
        "Source: lavita/MedQuAD on Hugging Face — 47,457 question-answer pairs",
        "Origin: NIH websites (GARD, GHR, NIDDK, NINDS, NCI, NHLBI, etc.)",
        "Chunking: 300–500 tokens per chunk, 50-token overlap for context continuity",
        "Embeddings: all-MiniLM-L6-v2 produces 384-dimensional dense vectors",
        "Index target: ~1,000–5,000 chunks stored in ChromaDB with full metadata",
        "Optional extension: WHO fact sheets and public drug leaflets (PDF → text)",
    ]
    _bullets(slide, Inches(0.55), Inches(4.1), Inches(12.2), Inches(2.8), details, sz=13)
    _footer(slide)


def slide_tech_stack(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Technology Stack")
    _img(slide, ASSETS / "tech_stack.png", Inches(0.5), Inches(1.35), Inches(8.5))

    tools = [
        ("Python 3.10+", "Core language"),
        ("Ollama", "Local LLM runtime"),
        ("Llama 3.2 / Mistral", "Generation model"),
        ("sentence-transformers", "Embedding model"),
        ("ChromaDB", "Vector database"),
        ("Streamlit", "Demo web UI"),
        ("rouge-score", "Answer evaluation"),
        ("pandas / numpy", "Data processing"),
    ]
    for i, (tool, role) in enumerate(tools):
        y = Inches(1.45 + i * 0.68)
        _rounded(slide, Inches(9.3), y, Inches(3.5), Inches(0.55), WHITE)
        _text(slide, Inches(9.45), y + Inches(0.05), Inches(1.8), Inches(0.45),
              tool, sz=12, bold=True, color=NAVY)
        _text(slide, Inches(11.3), y + Inches(0.1), Inches(1.5), Inches(0.4),
              role, sz=10, color=GRAY)
    _footer(slide)


def slide_implementation(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Implementation Plan", "Four structured development phases")

    phases = [
        (TEAL, "Phase 1 — Data Pipeline", [
            "Download MedQuAD from Hugging Face datasets API",
            "Clean HTML tags, normalize whitespace, deduplicate entries",
            "Split documents into overlapping chunks (300–500 tokens)",
            "Generate embeddings and build ChromaDB persistent index",
            "Deliverable: indexed corpus + ingestion statistics report",
        ]),
        (BLUE, "Phase 2 — RAG Core", [
            "Implement semantic search with configurable Top-K (default K=5)",
            "Design system prompt: answer only from provided context",
            "Integrate Ollama API for local text generation",
            "Add similarity score threshold to filter low-quality retrievals",
            "Deliverable: working rag.py CLI that answers questions",
        ]),
        (PURPLE, "Phase 3 — Demo UI", [
            "Build Streamlit interface with question input and history",
            "Display answer, source passages, and cosine similarity scores",
            "Show real-time latency counter and retrieval debug panel",
            "Add medical disclaimer banner and example question buttons",
            "Deliverable: browser-based live demo ready for presentation",
        ]),
        (ORANGE, "Phase 4 — Evaluation", [
            "Hold out 100 MedQuAD pairs as test set (never seen during indexing)",
            "Compute Recall@5, MRR, ROUGE-L, and latency distributions",
            "Generate bar charts and result tables for technical report",
            "Conduct 20-sample manual faithfulness audit",
            "Deliverable: results/ folder with CSV, plots, and summary",
        ]),
    ]
    for i, (col, title, items) in enumerate(phases):
        x = Inches(0.4 + (i % 2) * 6.5)
        y = Inches(1.4 + (i // 2) * 2.95)
        _rounded(slide, x, y, Inches(6.2), Inches(2.7), WHITE)
        _rect(slide, x, y, Inches(6.2), Inches(0.5), col)
        _text(slide, x + Inches(0.2), y + Inches(0.08), Inches(5.8), Inches(0.4),
              title, sz=13, bold=True, color=WHITE)
        _bullets(slide, x + Inches(0.2), y + Inches(0.6), Inches(5.8), Inches(2.0),
                 [f"• {it}" for it in items], sz=10, color=NAVY, spacing=3)
    _footer(slide)


def slide_prompt(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Prompt Engineering Strategy",
                "Constraining the LLM to grounded, cited responses")

    _img(slide, ASSETS / "ai_brain.jpg", Inches(0.55), Inches(1.45), Inches(4.8), Inches(5.2))

    prompt_text = (
        "SYSTEM:\n"
        "You are MedAssist, a medical information assistant.\n"
        "Answer ONLY using the provided context passages.\n"
        "If the context does not contain the answer, say\n"
        "\"I don't have enough information in my sources.\"\n"
        "Always cite which source passage you used.\n"
        "This is informational only — not medical advice.\n\n"
        "CONTEXT:\n"
        "[Source 1] (score: 0.87) Aspirin may cause...\n"
        "[Source 2] (score: 0.82) Common side effects...\n\n"
        "QUESTION: What are side effects of aspirin?\n\n"
        "ANSWER:"
    )
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(5.6), Inches(1.55), Inches(7.2), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    box.line.color.rgb = TEAL
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = prompt_text
    p.font.size = Pt(12)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xA8, 0xE6, 0xCF)

    rules = [
        "Temperature: 0.1 (low creativity, high factual consistency)",
        "Max tokens: 512 (concise medical answers)",
        "Top-K retrieval: 5 chunks (balance context vs. noise)",
        "Similarity threshold: 0.5 (reject irrelevant retrievals)",
    ]
    _bullets(slide, Inches(5.6), Inches(6.2), Inches(7.2), Inches(1.0), rules, sz=12, color=GRAY)
    _footer(slide)


def slide_metrics(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Evaluation Metrics",
                "Numerical performance measurement as required by the project brief")
    _img(slide, ASSETS / "metrics_dashboard.png", Inches(0.5), Inches(1.35), Inches(12.3))

    formulas = [
        "Recall@5 = |relevant ∩ top-5| / |relevant|",
        "MRR = (1/N) Σ 1/rank_i",
        "ROUGE-L = LCS-based F-measure vs. reference answer",
        "Latency = wall-clock time from query submit to full response",
    ]
    for i, f in enumerate(formulas):
        _rounded(slide, Inches(0.55 + i * 3.15), Inches(5.0), Inches(3.0), Inches(0.55),
                 WHITE, f, sz=9, tc=NAVY)
    _footer(slide)


def slide_demo(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Live Demo Plan", "June 17, 2026 — University Presentation")

    _img(slide, ASSETS / "data_analytics.jpg", Inches(0.55), Inches(1.45), Inches(5.0), Inches(5.2))

    demos = [
        ("Q1", "What are common side effects of aspirin?",
         "Shows drug information retrieval from MedQuAD"),
        ("Q2", "What is hypertension and how is it diagnosed?",
         "Shows multi-aspect medical condition explanation"),
        ("Q3", "What are the symptoms of diabetes?",
         "Shows symptom list with NIH source citations"),
        ("Q4", "Tell me about quantum physics.",
         "Failure case — system refuses out-of-domain questions"),
    ]
    for i, (q, question, note) in enumerate(demos):
        y = Inches(1.5 + i * 1.3)
        _rounded(slide, Inches(5.8), y, Inches(7.0), Inches(1.1), WHITE)
        _rounded(slide, Inches(5.95), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 TEAL if i < 3 else ORANGE, q, sz=12, bold=True)
        _text(slide, Inches(6.6), y + Inches(0.12), Inches(6.0), Inches(0.4),
              question, sz=13, bold=True, color=NAVY)
        _text(slide, Inches(6.6), y + Inches(0.55), Inches(6.0), Inches(0.4),
              note, sz=11, color=GRAY)

    _rounded(slide, Inches(5.8), Inches(6.75), Inches(7.0), Inches(0.42),
             TEAL, "Each demo shows: Answer  •  Source chunks  •  Similarity scores  •  Response time (ms)",
             sz=10, bold=True)
    _footer(slide)


def slide_timeline(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Project Timeline")
    _img(slide, ASSETS / "timeline.png", Inches(0.4), Inches(1.5), Inches(12.5))

    milestones = [
        "✓  Project plan & presentation submitted",
        "○  Week 1: Data pipeline complete",
        "○  Week 2: RAG + Streamlit UI working",
        "○  Week 3: Evaluation results + report draft",
        "★  June 17: Live demo + final report submission",
    ]
    _bullets(slide, Inches(0.55), Inches(4.8), Inches(12), Inches(2.2), milestones, sz=14)
    _footer(slide)


def slide_risks(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Risks & Mitigations")

    risks = [
        (ORANGE, "Slow CPU inference",
         "Ollama responses take >10s on weak hardware",
         "Use llama3.2:1b or mistral:7b-q4; limit max_tokens to 256"),
        (RGBColor(0xE7, 0x4C, 0x3C), "Low retrieval accuracy",
         "Wrong chunks retrieved → incorrect answers",
         "Tune chunk size (256/512), increase overlap, try bge-small-en"),
        (PURPLE, "LLM hallucination",
         "Model invents facts not in retrieved context",
         "Strict system prompt + temperature 0.1 + faithfulness audit"),
        (BLUE, "Ethical / legal concerns",
         "Users may treat output as medical diagnosis",
         "Prominent disclaimer; scope limited to educational Q&A"),
        (TEAL, "Dataset bias",
         "MedQuAD is English-only, US-centric NIH content",
         "Document limitation in report; note future multilingual extension"),
    ]
    for i, (col, risk, problem, fix) in enumerate(risks):
        y = Inches(1.4 + i * 1.15)
        _rounded(slide, Inches(0.55), y, Inches(12.2), Inches(1.0), WHITE)
        _rect(slide, Inches(0.55), y, Inches(0.12), Inches(1.0), col)
        _text(slide, Inches(0.85), y + Inches(0.1), Inches(2.5), Inches(0.35),
              risk, sz=13, bold=True, color=col)
        _text(slide, Inches(3.5), y + Inches(0.1), Inches(4.5), Inches(0.35),
              f"Risk: {problem}", sz=11, color=GRAY)
        _text(slide, Inches(8.2), y + Inches(0.1), Inches(4.5), Inches(0.8),
              f"→ {fix}", sz=11, color=NAVY)
    _footer(slide)


def slide_repo(prs):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), W, H, LIGHT)
    _header_bar(slide, "Repository Structure",
                "Organized codebase for reproducibility")

    structure = (
        "ML-pro/\n"
        "├── data/\n"
        "│   ├── raw/              # MedQuAD download\n"
        "│   └── processed/        # chunked JSONL files\n"
        "├── src/\n"
        "│   ├── ingest.py         # load & chunk MedQuAD\n"
        "│   ├── embed_index.py    # build ChromaDB index\n"
        "│   ├── rag.py            # retrieve + generate\n"
        "│   └── evaluate.py       # Recall@k, MRR, ROUGE-L\n"
        "├── app/\n"
        "│   └── streamlit_app.py  # live demo UI\n"
        "├── results/\n"
        "│   ├── metrics.json      # evaluation output\n"
        "│   └── plots/            # charts for report\n"
        "├── report/               # presentation & report\n"
        "├── requirements.txt\n"
        "└── README.md"
    )
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.55), Inches(1.45), Inches(6.0), Inches(5.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    box.line.color.rgb = TEAL
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = structure
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xA8, 0xE6, 0xCF)

    deliverables = [
        ("Dataset Sourcing", "MedQuAD + optional PDFs, with preprocessing logs"),
        ("Real-Time Demo", "Streamlit app with live Ollama inference"),
        ("Technical Report", "Architecture, training logs, evaluation tables"),
        ("Numerical Eval", "Recall@5, MRR, ROUGE-L, latency benchmarks"),
        ("Source Code", "Full reproducible pipeline on GitHub"),
    ]
    for i, (title, desc) in enumerate(deliverables):
        y = Inches(1.55 + i * 1.05)
        _rounded(slide, Inches(6.9), y, Inches(5.9), Inches(0.85), WHITE)
        _rounded(slide, Inches(7.05), y + Inches(0.15), Inches(0.4), Inches(0.4),
                 TEAL, str(i + 1), sz=11, bold=True)
        _text(slide, Inches(7.6), y + Inches(0.1), Inches(4.9), Inches(0.35),
              title, sz=13, bold=True, color=NAVY)
        _text(slide, Inches(7.6), y + Inches(0.48), Inches(4.9), Inches(0.35),
              desc, sz=11, color=GRAY)
    _footer(slide)


def slide_conclusion(prs):
    slide = _blank(prs)
    if _img(slide, ASSETS / "privacy_lock.jpg", Inches(0), Inches(0), W, H):
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = NAVY
        overlay.fill.transparency = 0.5
        overlay.line.fill.background()

    _text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
          "Conclusion", sz=40, bold=True, color=WHITE)

    points = [
        "MedAssist Local delivers private, source-grounded healthcare Q&A",
        "Local LLM + RAG = no cloud dependency, full data privacy",
        "MedQuAD provides a rich, NIH-backed medical knowledge base",
        "Rigorous numerical evaluation meets all project requirements",
        "Live Streamlit demo ready for June 17, 2026 presentation",
    ]
    for i, pt in enumerate(points):
        _text(slide, Inches(1.0), Inches(2.6 + i * 0.65), Inches(11), Inches(0.5),
              f"✓  {pt}", sz=17, color=WHITE)

    _rounded(slide, Inches(0.8), Inches(5.8), Inches(5), Inches(0.55),
             TEAL, "Thank You  —  Questions?", sz=20, bold=True)
    _text(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4),
          "[Your Name]  •  [your.email@university.edu]  •  GitHub: [your-repo]",
          sz=13, color=TEAL)


# ── main ───────────────────────────────────────────────────────────────────
def main() -> None:
    import subprocess
    import sys
    subprocess.run([sys.executable, str(ASSETS / "generate_assets.py")], check=True)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_agenda(prs)

    _section_slide(prs, "PART 1", "Background & Motivation",
                   "Understanding the problem and defining project scope", "doctor_consult.jpg")
    slide_problem(prs)
    slide_scope(prs)

    _section_slide(prs, "PART 2", "Technical Approach",
                   "RAG architecture, local LLM, and system design", "ai_brain.jpg")
    slide_why_rag(prs)
    slide_architecture(prs)
    slide_prompt(prs)
    slide_tech_stack(prs)

    _section_slide(prs, "PART 3", "Data & Implementation",
                   "Dataset sourcing, preprocessing, and development plan", "data_analytics.jpg")
    slide_dataset(prs)
    slide_implementation(prs)
    slide_repo(prs)

    _section_slide(prs, "PART 4", "Evaluation & Delivery",
                   "Metrics, demo plan, timeline, and risk management", "team_work.jpg")
    slide_metrics(prs)
    slide_demo(prs)
    slide_timeline(prs)
    slide_risks(prs)
    slide_conclusion(prs)

    out = OUTPUT
    try:
        prs.save(out)
    except PermissionError:
        out = HERE / "MedAssist_Local_Premium_v2.pptx"
        prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
